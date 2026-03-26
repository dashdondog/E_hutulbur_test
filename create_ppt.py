from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK = RGBColor(30, 41, 59)       # slate-800
MED = RGBColor(71, 85, 105)       # slate-500
LIGHT = RGBColor(148, 163, 184)   # slate-400
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(59, 130, 246)
BLUE_DARK = RGBColor(29, 78, 216)
PURPLE = RGBColor(139, 92, 246)
GREEN = RGBColor(16, 185, 129)
ORANGE = RGBColor(249, 115, 22)
RED = RGBColor(239, 68, 68)
BG_LIGHT = RGBColor(248, 250, 252)
BG_BLUE = RGBColor(239, 246, 255)

def add_bg(slide, color=BG_LIGHT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(8)
        p.level = 0
    return txBox

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, WHITE)

# Blue accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

# Main title area
add_shape_bg(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5), BG_BLUE, radius=0.05)

add_text_box(slide, Inches(2.2), Inches(1.8), Inches(8.8), Inches(1.2),
    "11-р ангийн хичээлийн систем", 44, DARK, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.2), Inches(3.0), Inches(8.8), Inches(0.8),
    "AI-д суурилсан хөтөлбөр & тест бэлтгэх платформ", 24, MED, False, PP_ALIGN.CENTER)

# Tech badges
badges = ["Next.js 16", "React 19", "Claude AI", "TypeScript", "Tailwind CSS"]
start_x = 3.0
for i, badge in enumerate(badges):
    shape = add_shape_bg(slide, Inches(start_x + i * 1.6), Inches(4.2), Inches(1.4), Inches(0.45), BLUE, radius=0.15)
    shape.text_frame.paragraphs[0].text = badge
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(2.2), Inches(5.5), Inches(8.8), Inches(0.5),
    "2026 он", 16, LIGHT, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Асуудал
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), RED)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.8), "Асуудал", 36, RED, True)

problems = [
    ("⏰", "Цаг хугацаа", "Багш нар хөтөлбөр, тест бэлтгэхэд маш их цаг зарцуулдаг. Нэг хичээлийн тест бэлтгэхэд дунджаар 3-4 цаг зарцуулна."),
    ("📚", "Гар ажиллагаа", "Сурах бичгийн агуулгыг гараар задлан шинжилж, сэдэв сэдвээр ялгах хүнд ажил. Алдаа гарах магадлал өндөр."),
    ("🔄", "Давтагдсан ажил", "Сэдэв бүрээр тест бэлдэх нь давтагдсан, нэг хэвийн ажил. 16 хичээл x олон сэдэв = маш их ажил."),
]

for i, (icon, title, desc) in enumerate(problems):
    y = 1.8 + i * 1.7
    add_shape_bg(slide, Inches(1), Inches(y), Inches(11.3), Inches(1.4), RGBColor(254, 242, 242), radius=0.03)
    add_text_box(slide, Inches(1.3), Inches(y + 0.15), Inches(0.6), Inches(0.6), icon, 32, DARK, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.1), Inches(y + 0.1), Inches(9.5), Inches(0.5), title, 22, RED, True)
    add_text_box(slide, Inches(2.1), Inches(y + 0.6), Inches(9.5), Inches(0.7), desc, 16, MED)

# ============================================================
# SLIDE 3: Шийдэл
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.8), "Шийдэл", 36, GREEN, True)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
    "Сурах бичгийн PDF-ийг AI-р уншуулж бүгдийг автоматжуулах", 22, DARK, False)

solutions = [
    ("📋", "Хөтөлбөр үүсгэх", "32 долоо хоногийн хөтөлбөр автомат үүсгэнэ", BLUE),
    ("📑", "Сэдвүүд гаргах", "Бүлэг, сэдэв, дэд сэдвүүдийг ялгаж гаргана", PURPLE),
    ("✅", "Тест бэлтгэх", "Сэдэв бүрээр олон төрлийн тест автомат үүсгэнэ", GREEN),
    ("📊", "Дүн гаргах", "Тест өгсний дараа шууд дүн, зөв хариулт харуулна", ORANGE),
]

for i, (icon, title, desc, color) in enumerate(solutions):
    x = 1 + i * 3.0
    add_shape_bg(slide, Inches(x), Inches(2.8), Inches(2.7), Inches(3.5), BG_LIGHT, radius=0.03)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.85), Inches(3.1), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = icon
    circle.text_frame.paragraphs[0].font.size = Pt(28)
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE

    add_text_box(slide, Inches(x + 0.1), Inches(4.3), Inches(2.5), Inches(0.5), title, 18, DARK, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.1), Inches(4.9), Inches(2.5), Inches(1), desc, 14, MED, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: Технологийн стек
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.8), "Технологийн стек", 36, DARK, True)

categories = [
    ("Frontend", BLUE, ["Next.js 16", "React 19", "Tailwind CSS 4", "TypeScript 5"]),
    ("Backend", PURPLE, ["Node.js API Routes", "File System Storage", "PDF Parse + OCR", "LocalStorage"]),
    ("AI", GREEN, ["Claude AI (Anthropic)", "Claude Sonnet 4", "JSON Structured Output", "Prompt Engineering"]),
    ("DevTools", ORANGE, ["ESLint", "PostCSS", "Tesseract OCR", "Poppler (pdftotext)"]),
]

for i, (cat_name, color, items) in enumerate(categories):
    x = 0.8 + i * 3.1
    # Header
    add_shape_bg(slide, Inches(x), Inches(1.8), Inches(2.8), Inches(0.6), color, radius=0.05)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].text = cat_name
    shape.text_frame.paragraphs[0].font.size = Pt(18)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Items
    add_shape_bg(slide, Inches(x), Inches(2.5), Inches(2.8), Inches(3.5), BG_LIGHT, radius=0.03)
    for j, item in enumerate(items):
        add_text_box(slide, Inches(x + 0.2), Inches(2.7 + j * 0.8), Inches(2.4), Inches(0.7),
            f"• {item}", 15, DARK)

# ============================================================
# SLIDE 5: Системийн бүтэц
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.8), "Системийн бүтэц", 36, DARK, True)

# Architecture boxes
arch_items = [
    (1.0, 3.0, "Хэрэглэгч\n(Browser)", BLUE),
    (3.8, 3.0, "Next.js\nFrontend", PURPLE),
    (6.6, 3.0, "API Routes\n(Backend)", GREEN),
    (9.4, 3.0, "Claude AI\n(Anthropic)", ORANGE),
]

for x, y, text, color in arch_items:
    shape = add_shape_bg(slide, Inches(x), Inches(y), Inches(2.2), Inches(1.3), color, radius=0.05)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

# Arrows between boxes
for x in [3.2, 6.0, 8.8]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3.4), Inches(0.6), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = LIGHT
    arrow.line.fill.background()

# Bottom layer
bottom_items = [
    (1.5, 5.0, "PDF Upload\n& Parse", RGBColor(59, 130, 246)),
    (4.5, 5.0, "Text Extraction\n(OCR)", RGBColor(139, 92, 246)),
    (7.5, 5.0, "AI Content\nGeneration", RGBColor(16, 185, 129)),
    (10.0, 5.0, "LocalStorage\n& Files", RGBColor(249, 115, 22)),
]

for x, y, text, color in bottom_items:
    shape = add_shape_bg(slide, Inches(x), Inches(y), Inches(2.2), Inches(1.0), color, radius=0.05)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 6: Хичээлүүд
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.8), "16 хичээл", 36, DARK, True)
add_text_box(slide, Inches(1), Inches(1.2), Inches(8), Inches(0.5),
    "11-р ангийн бүх үндсэн хичээлүүдийг хамарсан", 18, MED)

subjects_data = [
    ("📝", "Монгол хэл", "#3B82F6"), ("🖋️", "Монгол бичиг", "#8B5CF6"),
    ("📖", "Уран зохиол", "#EC4899"), ("📐", "Математик", "#EF4444"),
    ("⚛️", "Физик", "#F97316"), ("🧪", "Хими", "#10B981"),
    ("🧬", "Биологи", "#06B6D4"), ("🌍", "Газарзүй", "#14B8A6"),
    ("🏛️", "Түүх", "#A855F7"), ("⚖️", "Нийгмийн ухаан", "#6366F1"),
    ("🇬🇧", "Англи хэл", "#0EA5E9"), ("💻", "Мэдээлэл зүй", "#22C55E"),
    ("🔧", "Технологи", "#F59E0B"), ("🎨", "Дүрслэх урлаг", "#E11D48"),
    ("🎵", "Хөгжим", "#7C3AED"), ("⚽", "Биеийн тамир", "#059669"),
]

for i, (icon, name, hex_color) in enumerate(subjects_data):
    row = i // 4
    col = i % 4
    x = 1.0 + col * 3.0
    y = 2.0 + row * 1.2
    r, g, b = int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)

    add_shape_bg(slide, Inches(x), Inches(y), Inches(2.7), Inches(0.9), BG_LIGHT, radius=0.03)
    add_text_box(slide, Inches(x + 0.15), Inches(y + 0.15), Inches(0.5), Inches(0.6), icon, 24, DARK, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.7), Inches(y + 0.2), Inches(1.8), Inches(0.5), name, 16, RGBColor(r, g, b), True)

# ============================================================
# SLIDE 7: Файл байршуулах
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(8), Inches(0.8), "Файл байршуулах", 36, DARK, True)

# Upload zone mockup
add_shape_bg(slide, Inches(1), Inches(1.8), Inches(5.5), Inches(3.0), BG_BLUE, radius=0.03)
add_text_box(slide, Inches(1.5), Inches(2.2), Inches(4.5), Inches(0.6), "📂", 48, DARK, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(4.5), Inches(0.5),
    "PDF, TXT файлаа энд чирж оруулна уу", 16, MED, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.5), Inches(4.5), Inches(0.5),
    "эсвэл энд дарж файл сонгоно уу", 13, LIGHT, False, PP_ALIGN.CENTER)

# Features
features = [
    ("📕 PDF файл", "Сурах бичгийн PDF-ийг автомат уншина"),
    ("📄 TXT/MD файл", "Текст файлуудыг шууд ашиглана"),
    ("🔍 OCR дэмжлэг", "Зурган PDF-ийг Tesseract-аар уншина"),
    ("🗑️ Файл устгах", "Шаардлагагүй файлыг устгах боломжтой"),
]

for i, (title, desc) in enumerate(features):
    y = 1.8 + i * 1.1
    add_shape_bg(slide, Inches(7), Inches(y), Inches(5.5), Inches(0.9), BG_LIGHT, radius=0.03)
    add_text_box(slide, Inches(7.3), Inches(y + 0.05), Inches(5), Inches(0.4), title, 16, DARK, True)
    add_text_box(slide, Inches(7.3), Inches(y + 0.45), Inches(5), Inches(0.4), desc, 13, MED)

# ============================================================
# SLIDE 8: AI хөтөлбөр үүсгэх
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "AI хөтөлбөр үүсгэх", 36, DARK, True)

curriculum_parts = [
    ("📌 Нэр", "Хөтөлбөрийн нэр автомат тодорхойлно"),
    ("🎯 Зорилго", "3-5 өгүүлбэрээр зорилго тодорхойлно"),
    ("📝 Агуулга", "Долоо хоног бүрийн хуваарьтай дэлгэрэнгүй агуулга"),
    ("📏 Үнэлгээ", "Үнэлгээний шалгуурууд тодорхойлно"),
    ("📅 Хугацаа", "32 долоо хоногийн хөтөлбөр"),
]

for i, (icon_title, desc) in enumerate(curriculum_parts):
    y = 1.8 + i * 1.05
    add_shape_bg(slide, Inches(1), Inches(y), Inches(11.3), Inches(0.85), BG_LIGHT, radius=0.03)
    add_text_box(slide, Inches(1.3), Inches(y + 0.1), Inches(3), Inches(0.35), icon_title, 18, BLUE, True)
    add_text_box(slide, Inches(4.5), Inches(y + 0.1), Inches(7.5), Inches(0.6), desc, 16, MED)

add_shape_bg(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.7), BLUE, radius=0.1)
shape = slide.shapes[-1]
shape.text_frame.paragraphs[0].text = "📋 Хөтөлбөр үүсгэх  →  Claude Sonnet 4"
shape.text_frame.paragraphs[0].font.size = Pt(16)
shape.text_frame.paragraphs[0].font.color.rgb = WHITE
shape.text_frame.paragraphs[0].font.bold = True
shape.text_frame.paragraphs[0].font.name = "Arial"
shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 9: AI сэдэв гаргах
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "AI сэдэв гаргах", 36, DARK, True)
add_text_box(slide, Inches(1), Inches(1.3), Inches(10), Inches(0.5),
    "Сурах бичгээс бүх бүлэг/сэдвийг дэд сэдвүүдтэй нь автомат ялгаж гаргана", 18, MED)

# Example topics
example_topics = [
    ("1-р бүлэг. Кинематик", ["Шулуун шугаман хөдөлгөөн", "Чөлөөт уналт", "Хурдатгал"]),
    ("2-р бүлэг. Динамик", ["Ньютоны хуулиуд", "Үрэлтийн хүч", "Инерц"]),
    ("3-р бүлэг. Энерги", ["Кинетик энерги", "Потенциал энерги", "Энергийн хадгалалт"]),
]

for i, (topic_name, subtopics) in enumerate(example_topics):
    x = 1 + i * 4.0
    add_shape_bg(slide, Inches(x), Inches(2.3), Inches(3.7), Inches(4.0), BG_LIGHT, radius=0.03)

    add_shape_bg(slide, Inches(x), Inches(2.3), Inches(3.7), Inches(0.7), PURPLE, radius=0.03)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].text = topic_name
    shape.text_frame.paragraphs[0].font.size = Pt(15)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    for j, sub in enumerate(subtopics):
        add_shape_bg(slide, Inches(x + 0.2), Inches(3.3 + j * 0.8), Inches(3.3), Inches(0.6), WHITE, radius=0.03)
        add_text_box(slide, Inches(x + 0.4), Inches(3.35 + j * 0.8), Inches(2.9), Inches(0.5), f"• {sub}", 14, MED)

# ============================================================
# SLIDE 10: AI тест бэлтгэх
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "AI тест бэлтгэх", 36, DARK, True)

# Question types
q_types = [
    ("Олон сонголттой", "4 хариулт, 1 зөв\nA, B, C, D сонголт", BLUE, "ABCD"),
    ("Үнэн / Худал", "2 сонголт\nҮнэн эсвэл Худал", PURPLE, "Ү/Х"),
    ("Нээлттэй", "Чөлөөт хариулт\nБичвэр оруулах", GREEN, "..."),
]

for i, (title, desc, color, symbol) in enumerate(q_types):
    x = 1.2 + i * 4.0
    add_shape_bg(slide, Inches(x), Inches(1.8), Inches(3.5), Inches(2.2), BG_LIGHT, radius=0.03)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.25), Inches(2.0), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = symbol
    circle.text_frame.paragraphs[0].font.size = Pt(18)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(x + 0.2), Inches(3.1), Inches(3.1), Inches(0.4), title, 17, DARK, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.2), Inches(3.5), Inches(3.1), Inches(0.5), desc, 13, MED, False, PP_ALIGN.CENTER)

# Difficulty levels
add_text_box(slide, Inches(1), Inches(4.5), Inches(10), Inches(0.5), "Хүндийн зэрэг:", 20, DARK, True)

diff_levels = [
    ("Хялбар", "1 оноо", RGBColor(34, 197, 94)),
    ("Дунд", "2 оноо", RGBColor(234, 179, 8)),
    ("Хэцүү", "3 оноо", RGBColor(239, 68, 68)),
]

for i, (level, points, color) in enumerate(diff_levels):
    x = 1 + i * 4.0
    shape = add_shape_bg(slide, Inches(x), Inches(5.2), Inches(3.5), Inches(0.8), color, radius=0.1)
    shape.text_frame.paragraphs[0].text = f"{level}  —  {points}"
    shape.text_frame.paragraphs[0].font.size = Pt(18)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 11: Тест өгөх горим
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "Тест өгөх горим", 36, DARK, True)

test_features = [
    ("⏱️", "Цаг хэмжигч", "40 минутын хугацаатай. Цаг дуусвал автомат дуусгана. Үлдсэн цагийг байнга харуулна.", BLUE),
    ("✋", "Хариулт сонгох", "Олон сонголт, Үнэн/Худал асуултуудад дарж хариулна. Нээлттэй асуултад бичнэ.", PURPLE),
    ("📊", "Автомат дүн", "Дуусмагц шууд дүн гаргана: зөв/буруу тоо, нийт оноо, хувь (%).", GREEN),
    ("✅", "Зөв хариулт", "Буруу хариулсан асуултын зөв хариултыг ногоон өнгөөр харуулна.", ORANGE),
]

for i, (icon, title, desc, color) in enumerate(test_features):
    y = 1.6 + i * 1.35
    add_shape_bg(slide, Inches(1), Inches(y), Inches(11.3), Inches(1.15), BG_LIGHT, radius=0.03)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), Inches(y + 0.15), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = icon
    circle.text_frame.paragraphs[0].font.size = Pt(22)
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.5), Inches(y + 0.05), Inches(9.5), Inches(0.4), title, 19, DARK, True)
    add_text_box(slide, Inches(2.5), Inches(y + 0.5), Inches(9.5), Inches(0.6), desc, 14, MED)

# ============================================================
# SLIDE 12: Өгөгдлийн бүтэц
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "Өгөгдлийн бүтэц", 36, DARK, True)

models = [
    ("Subject", BLUE, ["id: string", "name: string", "icon: string", "color: string"]),
    ("Curriculum", PURPLE, ["name: string", "goal: string", "content: string", "criteria: string", "weeks: number"]),
    ("Topic", GREEN, ["name: string", "description: string", "subtopics: string[]", "order: number"]),
    ("Test / Question", ORANGE, ["name: string", "duration: number", "questions: Question[]", "type: multiple | truefalse", "correctAnswer: number"]),
]

for i, (name, color, fields) in enumerate(models):
    x = 0.6 + i * 3.2
    # Header
    add_shape_bg(slide, Inches(x), Inches(1.8), Inches(2.9), Inches(0.6), color, radius=0.05)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].text = name
    shape.text_frame.paragraphs[0].font.size = Pt(16)
    shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Fields
    add_shape_bg(slide, Inches(x), Inches(2.5), Inches(2.9), Inches(len(fields) * 0.6 + 0.3), BG_LIGHT, radius=0.03)
    for j, field in enumerate(fields):
        add_text_box(slide, Inches(x + 0.2), Inches(2.65 + j * 0.6), Inches(2.5), Inches(0.5),
            field, 13, MED, False, PP_ALIGN.LEFT, "Courier New")

# ============================================================
# SLIDE 13: API Endpoints
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "API Endpoints", 36, DARK, True)

# Table header
add_shape_bg(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(0.7), DARK, radius=0.02)
headers = [("Method", 1.2, 2.0), ("Endpoint", 3.5, 4.5), ("Тайлбар", 7.0, 5.0)]
for text, x, w in headers:
    add_text_box(slide, Inches(x), Inches(1.85), Inches(w), Inches(0.5), text, 16, WHITE, True)

endpoints = [
    ("POST", "/api/upload", "PDF/TXT файл байршуулах", GREEN),
    ("GET", "/api/upload", "Хичээлийн файлуудыг авах", BLUE),
    ("DELETE", "/api/upload", "Файл устгах", RED),
    ("PATCH", "/api/upload", "Текст дахин extract хийх", ORANGE),
    ("POST", "/api/generate/curriculum", "AI хөтөлбөр үүсгэх", GREEN),
    ("POST", "/api/generate/topics", "AI сэдвүүд гаргах", GREEN),
    ("POST", "/api/generate/tests", "AI тест бэлтгэх", GREEN),
]

for i, (method, endpoint, desc, method_color) in enumerate(endpoints):
    y = 2.6 + i * 0.6
    bg_color = BG_LIGHT if i % 2 == 0 else WHITE
    add_shape_bg(slide, Inches(1), Inches(y), Inches(11.3), Inches(0.55), bg_color, radius=0.01)

    # Method badge
    badge = add_shape_bg(slide, Inches(1.2), Inches(y + 0.08), Inches(1.1), Inches(0.38), method_color, radius=0.1)
    badge.text_frame.paragraphs[0].text = method
    badge.text_frame.paragraphs[0].font.size = Pt(11)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.name = "Courier New"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(3.5), Inches(y + 0.05), Inches(3.5), Inches(0.45), endpoint, 14, DARK, False, PP_ALIGN.LEFT, "Courier New")
    add_text_box(slide, Inches(7.0), Inches(y + 0.05), Inches(5), Inches(0.45), desc, 14, MED)

# ============================================================
# SLIDE 14: Ажиллах урсгал
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "Ажиллах урсгал", 36, DARK, True)

steps = [
    ("1", "Хичээл сонгох", "16 хичээлээс сонгоно", BLUE),
    ("2", "PDF байршуулах", "Сурах бичгээ оруулна", PURPLE),
    ("3", "AI үүсгэх", "Товч дарахад AI ажиллана", GREEN),
    ("4", "Хөтөлбөр", "32 долоо хоногийн хөтөлбөр", ORANGE),
    ("5", "Сэдвүүд", "Бүлэг, дэд сэдвүүд", RED),
    ("6", "Тест", "Сэдэв бүрээр тест", BLUE_DARK),
]

for i, (num, title, desc, color) in enumerate(steps):
    if i < 3:
        x = 1.0 + i * 4.0
        y = 1.8
    else:
        x = 1.0 + (i - 3) * 4.0
        y = 4.2

    add_shape_bg(slide, Inches(x), Inches(y), Inches(3.5), Inches(1.8), BG_LIGHT, radius=0.03)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1.25), Inches(y + 0.2), Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].text = num
    circle.text_frame.paragraphs[0].font.size = Pt(28)
    circle.text_frame.paragraphs[0].font.color.rgb = WHITE
    circle.text_frame.paragraphs[0].font.bold = True
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.1), Inches(3.1), Inches(0.35), title, 17, DARK, True, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.2), Inches(y + 1.4), Inches(3.1), Inches(0.35), desc, 13, MED, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15: Цаашдын хөгжүүлэлт
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.8), "Цаашдын хөгжүүлэлт", 36, DARK, True)

future_items = [
    ("🗄️", "Өгөгдлийн сан", "MongoDB/Supabase руу шилжиж, өгөгдлийг найдвартай хадгалах", BLUE),
    ("🔐", "Нэвтрэх систем", "Багш, сурагчийн нэвтрэх (Authentication) систем нэмэх", PURPLE),
    ("📈", "Дүнгийн тайлан", "Сурагч бүрийн дүн, ахиц дэвшлийн тайлан гаргах", GREEN),
    ("👥", "Олон хэрэглэгч", "Олон багш, сурагч зэрэг ашиглах боломж", ORANGE),
    ("📱", "Мобайл хувилбар", "Гар утаснаас ашиглах responsive дизайн сайжруулах", RED),
]

for i, (icon, title, desc, color) in enumerate(future_items):
    y = 1.6 + i * 1.1
    add_shape_bg(slide, Inches(1), Inches(y), Inches(11.3), Inches(0.9), BG_LIGHT, radius=0.03)

    add_shape_bg(slide, Inches(1.2), Inches(y + 0.1), Inches(0.7), Inches(0.7), color, radius=0.1)
    shape = slide.shapes[-1]
    shape.text_frame.paragraphs[0].text = icon
    shape.text_frame.paragraphs[0].font.size = Pt(20)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(2.2), Inches(y + 0.05), Inches(3), Inches(0.4), title, 18, DARK, True)
    add_text_box(slide, Inches(2.2), Inches(y + 0.45), Inches(9.8), Inches(0.4), desc, 14, MED)

# ============================================================
# SLIDE 16: Баярлалаа
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

add_shape_bg(slide, Inches(2), Inches(1.5), Inches(9.3), Inches(4.5), BG_BLUE, radius=0.05)

add_text_box(slide, Inches(2.5), Inches(2.0), Inches(8.3), Inches(1.2),
    "Анхаарал тавьсанд баярлалаа!", 44, DARK, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.5), Inches(3.3), Inches(8.3), Inches(0.8),
    "11-р ангийн хичээлийн систем", 28, BLUE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.5), Inches(4.2), Inches(8.3), Inches(0.6),
    "AI-д суурилсан боловсролын платформ", 20, MED, False, PP_ALIGN.CENTER)

# Separator
add_shape_bg(slide, Inches(5.5), Inches(5.0), Inches(2.3), Inches(0.04), BLUE)

add_text_box(slide, Inches(2.5), Inches(5.3), Inches(8.3), Inches(0.5),
    "Асуулт байна уу?", 22, LIGHT, False, PP_ALIGN.CENTER)

# Save
output_path = "/Users/dashdondoggansaikhan/education-system/presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
